from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

def add_slide(title, content):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = title
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content[0] if content else ""
    
    for item in content[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1 if item.startswith("-") else 0
        if item.startswith("-"):
            p.text = item[1:].strip()
            
    return slide

# Slide 1: Title
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "BÁO CÁO BÀI TẬP LỚN CUỐI KỲ - LỚP N03"
subtitle.text = "Ứng dụng Quản lý Thư viện trên thiết bị di động\nNhóm thực hiện:\n1. Nguyễn Trung Hiếu (23011098)\n2. Trương Bùi Huy Hiếu (23010885)\n3. Nguyễn Văn Trọng (23010748)\nĐại học Phenikaa"

# Slide 2: Agenda
add_slide(
    "MỤC LỤC / AGENDA",
    [
        "Tổng quan về hệ thống và lộ trình phát triển:",
        "- Câu 1 & 2: User Stories và Phân tích thiết kế",
        "- Câu 3: Sơ đồ UML (Class & Sequence)",
        "- Câu 4 & 5: Thiết kế giao diện & Layout mẫu",
        "- Câu 6 & 7: Triển khai Code & Database ORM",
        "- Câu 8: Kiểm thử & Unit Test",
        "- Câu 9 & 10: Tổ chức Nhóm & Git Commits"
    ]
)

# Slide 3: Câu 1
add_slide(
    "CÂU 1: USER STORIES (CÂU CHUYỆN NGƯỜI DÙNG)",
    [
        "Hệ thống thiết kế xoay quanh vai trò Thủ thư (Librarian):",
        "- Là thủ thư, tôi muốn đăng nhập an toàn vào hệ thống quản lý.",
        "- Là thủ thư, tôi muốn xem Dashboard thống kê để quản lý sách trực quan.",
        "- Là thủ thư, tôi muốn Thêm/Sửa/Xóa chi tiết một cuốn sách.",
        "- Là thủ thư, tôi muốn tạo nhanh Phiếu Mượn/Trả sách cho độc giả.",
        "- Là thủ thư, tôi muốn tự động hóa việc cập nhật tồn kho khi sách được mượn/trả."
    ]
)

# Slide 4: Câu 2
add_slide(
    "CÂU 2: PHÂN TÍCH YÊU CẦU, ĐỐI TƯỢNG VÀ MỐI QUAN HỆ",
    [
        "Yêu cầu: App mượt mà, lưu trữ dữ liệu thời gian thực trên cloud.",
        "Các đối tượng chính (Entities):",
        "- 1. AppUser: Quản lý thông tin thủ thư.",
        "- 2. BookModel: Đối tượng sách (tên, tác giả, số lượng).",
        "- 3. LoanModel: Đối tượng phiếu giao dịch (Mượn hoặc Trả).",
        "Mối quan hệ:",
        "- 1 AppUser xử lý N LoanModel (Phiếu).",
        "- Mối quan hệ N-M giữa LoanModel và BookModel thông qua giỏ hàng (LoanBookItem)."
    ]
)

# Slide 5: Câu 3 (Class Diagram)
add_slide(
    "CÂU 3: SƠ ĐỒ CẤU TRÚC LỚP (CLASS DIAGRAM)",
    [
        "Cấu trúc lớp hỗ trợ Object Relational Mapping (ORM):",
        "- Lớp BookModel, LoanModel, AppUser đều có constructor từ Map.",
        "- Chứa các hàm toMap() và fromMap()/fromJson().",
        "- Các Extension (như LoanModelExtension) giúp tái sử dụng logic tính toán.",
        "(Giành không gian trống này để dán ảnh Sơ đồ Class Diagram)"
    ]
)

# Slide 6: Câu 3 (Sequence/Activity Diagrams)
add_slide(
    "CÂU 3: SƠ ĐỒ THUẬT TOÁN (SEQUENCE/ACTIVITY)",
    [
        "Nhóm đã xây dựng 5 sơ đồ luồng hoạt động cốt lõi:",
        "- 1. Sơ đồ đăng nhập và xác thực.",
        "- 2. Sơ đồ quá trình thêm/sửa sách.",
        "- 3. Sơ đồ tạo phiếu mượn (trừ tồn kho).",
        "- 4. Sơ đồ tạo phiếu trả (cộng dồn tồn kho).",
        "- 5. Sơ đồ thuật toán lọc dữ liệu Dashboard theo thời gian.",
        "(Chèn ảnh các sơ đồ vào khu vực này)"
    ]
)

# Slide 7: Câu 4
add_slide(
    "CÂU 4: THIẾT KẾ MÀN HÌNH & FLOW OF WORK",
    [
        "Wireframes / Mockup Screens:",
        "- Giao diện sử dụng Material 3, Clean UI (Trắng & Xanh Navy).",
        "Logic luồng bài toán (Flow of Work):",
        "- Màn hình Đăng nhập -> HomeScreen (Bottom Navigation Bar).",
        "- Bottom Nav chứa 5 Tab: Trang chủ, Sách, Dashboard, Mượn, Trả.",
        "- Nút Settings tích hợp gọn gàng trong góc của Dashboard."
    ]
)

# Slide 8: Câu 5
add_slide(
    "CÂU 5: THỰC HIỆN LAYOUT MẪU YÊU CẦU",
    [
        "Màn hình thiết kế theo mẫu IMG_5518.PNG do giảng viên yêu cầu.",
        "- Responsive với mọi kích thước điện thoại.",
        "- Phối màu tuân thủ 100% bản mẫu gốc.",
        "- Đảm bảo cả tỷ lệ Padding, Margin và Typorgraphy.",
        "(Chèn ảnh ghép 2 bên: Ảnh mẫu IMG_5518 vs Ảnh ứng dụng thực tế)"
    ]
)

# Slide 9: Câu 6
add_slide(
    "CÂU 6: THỰC HIỆN CODE THEO USER STORIES",
    [
        "Triển khai toàn diện tính năng ứng dụng thực tiễn:",
        "- Tích hợp Firebase Auth cho đăng nhập.",
        "- Màn hình Add/Edit book kết hợp upload hình ảnh Cloudinary.",
        "- Màn hình Tạo phiếu tự động tính toán tổng số sách.",
        "- Triển khai Syncfusion Charts cho Dashboard báo cáo thống kê."
    ]
)

# Slide 10: Câu 7
add_slide(
    "CÂU 7: KẾT NỐI DATABASE VÀ KỸ THUẬT ORM",
    [
        "Cơ sở dữ liệu đám mây Firebase Firestore (NoSQL Real-time).",
        "Kỹ thuật Object Relational Mapping (ORM):",
        "- DocumentSnapshot -> Map -> Model Object.",
        "- Không gọi trực tiếp field String/Int từ Firebase lên UI.",
        "- Toàn bộ UI widget như ListView, GridView đều render dựa trên danh sách các Object mạnh (Strong-typed Objects)."
    ]
)

# Slide 11: Câu 8
add_slide(
    "CÂU 8: KIỂM THỬ, BẮT LỖI VÀ UNIT TEST",
    [
        "Kiểm thử chức năng và bắt lỗi ngoại lệ:",
        "- Bắt lỗi Form (TextFormField validate).",
        "- Bắt lỗi logic (Mượn lố số lượng có sẵn trong kho).",
        "Đơn vị kiểm định (Unit Test):",
        "- Khởi tạo test directory, sử dụng thư viện flutter_test.",
        "- Hoàn thành xuất sắc 10/10 test cases cho toàn bộ Models."
    ]
)

# Slide 12: Câu 9
add_slide(
    "CÂU 9: BÁO CÁO, LINK DEMO VÀ PHÂN CÔNG",
    [
        "Tài nguyên Đồ án (Đính kèm trong báo cáo giấy):",
        "- Link Github: https://github.com/Accnoname/project-mobile-app-library",
        "- Link Demo: (Điền link Firebase Hosting hoặc Youtube Demo)",
        "Đóng góp của nhóm:",
        "- Nguyễn Trung Hiếu (23011098): Phân tích hệ thống, vẽ sơ đồ UML, viết báo cáo.",
        "- Trương Bùi Huy Hiếu (23010885): Code giao diện, Dashboard thống kê, viết Unit Test.",
        "- Nguyễn Văn Trọng (23010748): Thiết lập CSDL Firebase, ORM mapping, viết Unit Test."
    ]
)

# Slide 13: Câu 10
add_slide(
    "CÂU 10: LỊCH SỬ CODE VÀ GIT COMMITS",
    [
        "Kỹ năng làm việc nhóm qua công cụ quản lý mã nguồn Git:",
        "- Commit có cấu trúc rõ ràng (Convention: feat, fix, chore...).",
        "- Các nhánh độc lập (VD: nhánh main, nhánh giao-dien).",
        "- Tần suất commit phản ánh đúng tiến trình và công sức thực tế của từng sinh viên.",
        "(Chèn ảnh chụp màn hình tab Commits trên kho lưu trữ Github)"
    ]
)

# Slide 14: Demo
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "DEMO ỨNG DỤNG THỰC TẾ"
subtitle.text = "Xin trân trọng kính mời Hội đồng xem Demo ứng dụng\nvà đưa ra nhận xét / câu hỏi (Q&A)."

# Save presentation
file_path = "Bao_Cao_Do_An_N03_Update.pptx"
prs.save(file_path)
print(f"Presentation saved to {file_path}")
